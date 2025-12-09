import streamlit as st
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import json
import io

# પેજ સેટિંગ
st.set_page_config(page_title="Modi Sir's PPT Maker", page_icon="📚")

st.title("📚 Education PPT Maker (GSEB)")
st.write("Created by: **Mukesh Modi** (Govt. Sec. School, Vantdau)")

# ઇનપુટ વિભાગ
api_key = st.text_input("તમારી Gemini API Key અહીં નાખો:", type="password")
topic = st.text_input("ટોપિકનું નામ લખો (દા.ત. માનવ પાચનતંત્ર):")

# PPT બનાવવાનું ફંક્શન
def create_ppt(topic, slides_data):
    prs = Presentation()
    
    # બ્લેક થીમ ફંક્શન
    def set_black(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0,0,0)

    # 1. ટાઈટલ સ્લાઈડ
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_black(slide)
    
    t = slide.shapes.title
    t.text = topic
    t.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,0) # Yellow
    t.text_frame.paragraphs[0].font.bold = True
    
    s = slide.placeholders[1]
    s.text = "Created by Mukesh Modi\nGovt. Sec. And High. Sec. School, Vantdau"
    s.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255) # White

    # 2. કન્ટેન્ટ સ્લાઈડ્સ
    for item in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        set_black(slide)
        
        # ટાઈટલ
        t = slide.shapes.title
        t.text = item['title']
        t.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,0)
        t.text_frame.paragraphs[0].font.bold = True
        
        # કન્ટેન્ટ
        c = slide.placeholders[1]
        c.text = item['content']
        for p in c.text_frame.paragraphs:
            p.font.color.rgb = RGBColor(255,255,255)
            p.font.size = Pt(22)
            
    binary_output = io.BytesIO()
    prs.save(binary_output)
    binary_output.seek(0)
    return binary_output

# બટન ક્લિક થાય ત્યારે
if st.button("PPT બનાવો 🚀"):
    if api_key and topic:
        try:
            genai.configure(api_key=api_key)
            
            # --- સુધારો અહીં કર્યો છે (gemini-1.5-flash) ---
            model = genai.GenerativeModel('gemini-1.5-flash')
            
            with st.spinner('Gemini વિચારી રહ્યું છે...'):
                prompt = f"""
                Create 5 slides for Grade 10 Gujarati medium on '{topic}'.
                Strictly return JSON format:
                [
                    {{"title": "Slide Title", "content": "Point 1\\nPoint 2"}}
                ]
                No markdown, just JSON text.
                """
                
                res = model.generate_content(prompt)
                
                # ડેટા ક્લીનિંગ
                clean_text = res.text.replace("```json", "").replace("```", "")
                data = json.loads(clean_text)
                
                # PPT બનાવવી
                ppt_file = create_ppt(topic, data)
                
                st.success("PPT બની ગઈ!")
                
                # ડાઉનલોડ બટન
                st.download_button(
                    label="📥 ડાઉનલોડ PPT",
                    data=ppt_file,
                    file_name=f"{topic}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
        except Exception as e:
            st.error(f"ભૂલ આવી: {e}")
    else:
        st.warning("પહેલા API Key અને ટોપિક બંને નાખો.")
