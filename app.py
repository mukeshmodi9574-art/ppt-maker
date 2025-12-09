import streamlit as st
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import json
import io

st.set_page_config(page_title="Modi Sir's PPT Maker", page_icon="📚")

st.title("📚 Education PPT Maker")
st.write("Created by: **Mukesh Modi** (Govt. Sec. School, Vantdau)")

# API Key Input
api_key = st.text_input("તમારી Gemini API Key અહીં નાખો:", type="password")
topic = st.text_input("ટોપિકનું નામ લખો (દા.ત. પાચનતંત્ર):")

def create_ppt(topic, slides_data):
    prs = Presentation()
    def set_black(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0,0,0)

    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_black(slide)
    t = slide.shapes.title
    t.text = topic
    t.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,0)
    t.text_frame.paragraphs[0].font.bold = True
    s = slide.placeholders[1]
    s.text = "Created by Mukesh Modi"
    s.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)

    # Content Slides
    for item in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        set_black(slide)
        t = slide.shapes.title
        t.text = item['title']
        t.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,0)
        t.text_frame.paragraphs[0].font.bold = True
        c = slide.placeholders[1]
        c.text = item['content']
        for p in c.text_frame.paragraphs:
            p.font.color.rgb = RGBColor(255,255,255)
            p.font.size = Pt(22)
            
    binary_output = io.BytesIO()
    prs.save(binary_output)
    binary_output.seek(0)
    return binary_output

if st.button("PPT બનાવો"):
    if api_key and topic:
        try:
            genai.configure(api_key=api_key)
            model = genai.GenerativeModel('gemini-pro')
            with st.spinner('લખાઈ રહ્યું છે...'):
                prompt = f"Create 5 slides for Grade 10 Gujarati medium on '{topic}'. JSON format: [{{'title':'T','content':'C'}}]"
                res = model.generate_content(prompt)
                data = json.loads(res.text.replace("```json","").replace("```",""))
                ppt_file = create_ppt(topic, data)
                st.success("PPT બની ગઈ!")
                st.download_button("ડાઉનલોડ કરો", ppt_file, f"{topic}.pptx")
        except Exception as e:
            st.error(f"ભૂલ: {e}")
    else:
        st.warning("API Key અને ટોપિક બંને નાખો.")
